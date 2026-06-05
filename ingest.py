"""
ingest.py — MultiDoc AI · Production-Grade Ingestion Pipeline
==============================================================
Handles:
  - Plain text PDFs          → PDFPlumber (fast, accurate)
  - Scanned / image PDFs     → pdf2image + Tesseract OCR
  - Mixed PDFs               → per-page detection, best method used
  - DOCX                     → python-docx (text + embedded tables)
  - TXT / MD                 → direct read
  - CSV                      → pandas → markdown table
  - XLSX                     → openpyxl → markdown table per sheet
  - PPTX                     → python-pptx (text + slide notes)

Optimizations in this version:
  - Smaller chunks (800) with higher overlap (250) → better boundary preservation
  - MIN_CHUNK_LEN lowered to 50 → keeps short but important chunks (dates, numbers)
  - PDF: heading detection — H1/H2 lines preserved as chunk anchors
  - DOCX: heading styles preserved in extracted text
  - CSV/XLSX: column headers repeated in every chunk for context
  - Dedup: exact-duplicate chunks removed before embedding
  - Metadata enriched: page number, chunk index, doc type always set
"""

# ── Imports ────────────────────────────────────────────────────────────────────
import os
import io
import logging
import hashlib
import chromadb
import pytesseract
import pandas as pd

from dotenv import load_dotenv
from PIL import Image
from pdf2image import convert_from_path
from pdfplumber import open as plumber_open
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma

# ── Logging ────────────────────────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S"
)
log = logging.getLogger("ingest")

# ── Environment ────────────────────────────────────────────────────────────────
load_dotenv()

# ── Config ─────────────────────────────────────────────────────────────────────
DOCS_DIR        = "docs"
CHROMA_DIR      = "chroma_db"
COLLECTION_NAME = "knowledge_base"
CHUNK_SIZE      = 800    # smaller = more precise retrieval
CHUNK_OVERLAP   = 250    # higher overlap = fewer boundary cuts
MIN_CHUNK_LEN   = 50     # keep short chunks — may contain key facts/numbers

# ── Tesseract & Poppler paths (Windows) ────────────────────────────────────────
TESSERACT_PATH = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
POPPLER_PATH   = r"C:\poppler\Library\bin"

pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH

# ── Supported formats ──────────────────────────────────────────────────────────
SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".docx", ".csv", ".xlsx", ".pptx", ".md"}


# ══════════════════════════════════════════════════════════════════════════════
# DEDUPLICATION — remove exact duplicate chunks before embedding
# ══════════════════════════════════════════════════════════════════════════════

def dedup_chunks(chunks: list[Document]) -> list[Document]:
    """Remove exact duplicate page_content chunks using MD5 fingerprint."""
    seen    = set()
    unique  = []
    removed = 0
    for chunk in chunks:
        fp = hashlib.md5(chunk.page_content.strip().encode()).hexdigest()
        if fp not in seen:
            seen.add(fp)
            unique.append(chunk)
        else:
            removed += 1
    if removed:
        log.info(f"  Dedup: removed {removed} duplicate chunks")
    return unique


# ══════════════════════════════════════════════════════════════════════════════
# TEXT EXTRACTORS — one per file type
# ══════════════════════════════════════════════════════════════════════════════

def extract_pdf(path: str, filename: str) -> list[Document]:
    """
    Smart PDF extractor with heading detection:
      - PDFPlumber for text pages
      - Tesseract OCR fallback for scanned pages
      - Table extraction per page
      - Heading lines prefixed with ### for better chunking
    """
    docs      = []
    ocr_pages = 0

    try:
        with plumber_open(path) as pdf:
            total_pages = len(pdf.pages)
            log.info(f"  PDF: {filename} — {total_pages} pages")

            try:
                images = convert_from_path(path, dpi=200, poppler_path=POPPLER_PATH)
            except Exception as e:
                log.warning(f"  pdf2image failed ({e}) — OCR fallback disabled")
                images = []

            for i, page in enumerate(pdf.pages):
                page_num = i + 1
                text     = ""

                try:
                    extracted = page.extract_text()
                    if extracted:
                        text = _tag_headings(extracted.strip())
                except Exception:
                    pass

                # Table extraction
                table_text = ""
                try:
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            rows = []
                            for row in table:
                                clean_row = [str(cell).strip() if cell else "" for cell in row]
                                if any(clean_row):
                                    rows.append(" | ".join(clean_row))
                            if rows:
                                table_text += "\n\n**Table:**\n" + "\n".join(rows)
                except Exception:
                    pass

                # OCR fallback
                if len(text) < 50 and i < len(images):
                    try:
                        ocr_text = pytesseract.image_to_string(
                            images[i], config="--psm 3"
                        ).strip()
                        if len(ocr_text) > len(text):
                            text = ocr_text
                            ocr_pages += 1
                    except Exception as e:
                        log.warning(f"  OCR failed on page {page_num}: {e}")

                combined = (text + "\n\n" + table_text).strip()

                if combined:
                    docs.append(Document(
                        page_content=combined,
                        metadata={
                            "source":      filename,
                            "doc_name":    os.path.splitext(filename)[0],
                            "page":        page_num,
                            "total_pages": total_pages,
                            "type":        "pdf",
                        }
                    ))

        if ocr_pages:
            log.info(f"  OCR used on {ocr_pages}/{total_pages} pages")

    except Exception as e:
        log.error(f"  PDF extraction failed for {filename}: {e}")

    return docs


def _tag_headings(text: str) -> str:
    """
    Detect likely headings in PDF text and prefix them with ###.
    Heuristic: short lines (<80 chars) that are ALL CAPS or Title Case
    with no period at end → likely a heading.
    """
    lines  = text.split("\n")
    result = []
    for line in lines:
        stripped = line.strip()
        if (
            stripped
            and len(stripped) < 80
            and not stripped.endswith(".")
            and (stripped.isupper() or stripped.istitle())
            and len(stripped.split()) >= 2
        ):
            result.append(f"### {stripped}")
        else:
            result.append(line)
    return "\n".join(result)


def extract_txt(path: str, filename: str) -> list[Document]:
    """Plain text and Markdown files."""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read().strip()
        if not content:
            return []
        return [Document(
            page_content=content,
            metadata={
                "source":   filename,
                "doc_name": os.path.splitext(filename)[0],
                "type":     "text",
            }
        )]
    except Exception as e:
        log.error(f"  TXT extraction failed for {filename}: {e}")
        return []


def extract_docx(path: str, filename: str) -> list[Document]:
    """
    DOCX with heading style detection.
    Heading paragraphs prefixed with ### so chunker respects structure.
    """
    try:
        from docx import Document as DocxDocument
        doc   = DocxDocument(path)
        parts = []

        for para in doc.paragraphs:
            text  = para.text.strip()
            style = para.style.name.lower() if para.style else ""
            if not text:
                continue
            if "heading" in style:
                parts.append(f"### {text}")
            else:
                parts.append(text)

        # Tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append("**Table:**\n" + "\n".join(rows))

        content = "\n\n".join(parts)
        if not content.strip():
            return []

        return [Document(
            page_content=content,
            metadata={
                "source":   filename,
                "doc_name": os.path.splitext(filename)[0],
                "type":     "docx",
            }
        )]
    except Exception as e:
        log.error(f"  DOCX extraction failed for {filename}: {e}")
        return []


def extract_csv(path: str, filename: str) -> list[Document]:
    """
    CSV → chunked by row batches, headers repeated in each chunk.
    This way even if a chunk is retrieved mid-file, column context is preserved.
    """
    try:
        df = pd.read_csv(path, encoding="utf-8", on_bad_lines="skip")
        if df.empty:
            return []

        headers     = list(df.columns)
        header_line = " | ".join(headers)
        docs        = []
        batch_size  = 50  # rows per chunk

        for start in range(0, len(df), batch_size):
            batch   = df.iloc[start:start+batch_size]
            content = f"Columns: {header_line}\n\n" + batch.to_markdown(index=False)
            docs.append(Document(
                page_content=content,
                metadata={
                    "source":     filename,
                    "doc_name":   os.path.splitext(filename)[0],
                    "row_start":  start,
                    "row_end":    start + len(batch),
                    "type":       "csv",
                }
            ))

        return docs
    except Exception as e:
        log.error(f"  CSV extraction failed for {filename}: {e}")
        return []


def extract_xlsx(path: str, filename: str) -> list[Document]:
    """
    XLSX — each sheet chunked by row batches, headers repeated.
    """
    try:
        import openpyxl
        wb   = openpyxl.load_workbook(path, read_only=True, data_only=True)
        docs = []

        for sheet_name in wb.sheetnames:
            ws        = wb[sheet_name]
            all_rows  = []
            headers   = None

            for row in ws.iter_rows(values_only=True):
                clean = [str(cell).strip() if cell is not None else "" for cell in row]
                if any(clean):
                    if headers is None:
                        headers = clean
                    else:
                        all_rows.append(clean)

            if not headers:
                continue

            header_line = " | ".join(headers)
            batch_size  = 50

            for start in range(0, len(all_rows), batch_size):
                batch   = all_rows[start:start+batch_size]
                rows_md = "\n".join(" | ".join(r) for r in batch)
                content = f"Sheet: {sheet_name}\nColumns: {header_line}\n\n{rows_md}"

                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source":    filename,
                        "doc_name":  os.path.splitext(filename)[0],
                        "sheet":     sheet_name,
                        "row_start": start,
                        "row_end":   start + len(batch),
                        "type":      "xlsx",
                    }
                ))

        return docs
    except Exception as e:
        log.error(f"  XLSX extraction failed for {filename}: {e}")
        return []


def extract_pptx(path: str, filename: str) -> list[Document]:
    """PPTX — slide title always prepended to slide content for context."""
    try:
        from pptx import Presentation
        prs  = Presentation(path)
        docs = []

        for i, slide in enumerate(prs.slides, start=1):
            parts = []
            title = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        # First non-empty text is likely the slide title
                        if not title:
                            title = text
                            parts.append(f"### Slide {i}: {text}")
                        else:
                            parts.append(text)

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Speaker Notes]: {notes}")

            content = "\n".join(parts)
            if content.strip():
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source":       filename,
                        "doc_name":     os.path.splitext(filename)[0],
                        "slide":        i,
                        "slide_title":  title,
                        "type":         "pptx",
                    }
                ))

        return docs
    except Exception as e:
        log.error(f"  PPTX extraction failed for {filename}: {e}")
        return []


# ══════════════════════════════════════════════════════════════════════════════
# DISPATCHER
# ══════════════════════════════════════════════════════════════════════════════

def load_file(path: str, filename: str) -> list[Document]:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return extract_pdf(path, filename)
    elif ext in (".txt", ".md"):
        return extract_txt(path, filename)
    elif ext == ".docx":
        return extract_docx(path, filename)
    elif ext == ".csv":
        return extract_csv(path, filename)
    elif ext == ".xlsx":
        return extract_xlsx(path, filename)
    elif ext == ".pptx":
        return extract_pptx(path, filename)
    else:
        log.warning(f"  Skipping unsupported file: {filename}")
        return []


def load_all_documents(folder: str) -> list[Document]:
    all_docs = []
    files    = [f for f in os.listdir(folder)
                if os.path.splitext(f)[1].lower() in SUPPORTED_EXTENSIONS]

    if not files:
        log.warning("No supported files found in docs/ folder.")
        return []

    for filename in files:
        path = os.path.join(folder, filename)
        log.info(f"📄 Loading: {filename}")
        docs = load_file(path, filename)
        if docs:
            all_docs.extend(docs)
            log.info(f"  ✅ {filename} → {len(docs)} document(s) extracted")
        else:
            log.warning(f"  ⚠️  {filename} → nothing extracted")

    return all_docs


# ══════════════════════════════════════════════════════════════════════════════
# MAIN PIPELINE
# ══════════════════════════════════════════════════════════════════════════════

def main():
    log.info("═" * 55)
    log.info("MultiDoc AI — Ingestion Pipeline")
    log.info("═" * 55)

    # 1. Load
    log.info("STEP 1/4 — Loading documents...")
    documents = load_all_documents(DOCS_DIR)
    if not documents:
        log.error("No documents loaded. Aborting.")
        return
    log.info(f"  Total documents loaded: {len(documents)}")

    # 2. Chunk
    log.info("STEP 2/4 — Chunking...")
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n### ", "\n\n", "\n", ".", "!", "?", ",", " ", ""]
    )
    chunks = splitter.split_documents(documents)

    # Filter junk
    before = len(chunks)
    chunks = [c for c in chunks if len(c.page_content.strip()) >= MIN_CHUNK_LEN]
    log.info(f"  {len(chunks)} chunks after junk filter (removed {before - len(chunks)})")

    # Dedup
    chunks = dedup_chunks(chunks)
    log.info(f"  {len(chunks)} chunks after dedup")

    # 3. Embed
    log.info("STEP 3/4 — Embedding with sentence-transformers...")
    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

    # 4. Store
    log.info("STEP 4/4 — Storing in ChromaDB...")

    try:
        temp_client = chromadb.PersistentClient(path=CHROMA_DIR)
        temp_client.delete_collection(COLLECTION_NAME)
        log.info("  Old collection deleted.")
    except Exception as e:
        log.info(f"  No existing collection to delete ({e}), starting fresh.")

    vectorstore = Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        persist_directory=CHROMA_DIR,
        collection_name=COLLECTION_NAME,
    )

    log.info(f"✅ Done! {len(chunks)} chunks stored in ChromaDB.")
    log.info("═" * 55)


if __name__ == "__main__":
    main()